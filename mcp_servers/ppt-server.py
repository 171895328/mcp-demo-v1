import os
import sys
import io
from mcp.server.fastmcp import FastMCP
from typing import Optional, List, Dict, Any, Union, Tuple

# 标记库是否已安装
pptx_installed = True

# 尝试导入python-pptx库，如果没有安装则标记为未安装但不退出
try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError:
    print("警告: 未检测到python-pptx库，PowerPoint功能将不可用")
    print("请使用以下命令安装: pip install python-pptx")
    pptx_installed = False

# 尝试导入Pillow库，用于图片处理
pillow_installed = True
try:
    from PIL import Image
except ImportError:
    print("警告: 未检测到Pillow库，图片处理功能将受限")
    print("请使用以下命令安装: pip install Pillow")
    pillow_installed = False

# 创建一个MCP服务器，保持名称与配置文件一致
mcp = FastMCP("office editor")


@mcp.tool("open_powerpoint_presentation","打开一个现有的PowerPoint演示文稿并读取其基本信息。")
def open_powerpoint_presentation(file_path: str) -> str:
    """
    打开一个现有的PowerPoint演示文稿并读取其基本信息。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
    
    Returns:
        演示文稿的基本信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法打开PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片信息
        slide_count = len(prs.slides)
        
        # 获取每张幻灯片的基本信息
        slides_info = []
        for i, slide in enumerate(prs.slides):
            slide_title = "无标题"
            # 尝试获取幻灯片标题
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and (shape.name.startswith("Title") or "标题" in shape.name):
                    slide_title = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                    break
            
            # 计算幻灯片上的形状数量
            shape_count = len(slide.shapes)
            
            slides_info.append(f"幻灯片 {i+1}: {slide_title} (包含 {shape_count} 个形状)")
        
        # 构建演示文稿信息
        presentation_info = (
            f"文件名: {os.path.basename(file_path)}\n"
            f"幻灯片数量: {slide_count}\n\n"
            f"幻灯片概览:\n" + "\n".join(slides_info)
        )
        
        return presentation_info
    except Exception as e:
        return f"打开PowerPoint演示文稿时出错: {str(e)}"



if __name__ == "__main__":
    try:
        # 启动MCP服务器
        mcp.run()
    except KeyboardInterrupt:
        # 优雅地处理Ctrl+C中断
        print("服务器已停止")
    except Exception as e:
        # 处理其他异常
        print(f"服务器运行时出错: {str(e)}")
        sys.exit(1)